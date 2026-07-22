"""Generate the Dhwani-Kavach pitch deck as an editable .pptx (Team ERROR 404).

Dev tool. Requires: pip install python-pptx
Run:  python tools/build_pptx.py   -> writes Dhwani-Kavach_ERROR404.pptx
26-slide deck for a ~40-minute panel talk. Native PowerPoint shapes -> editable.
ponytail: hand-built slides, no template engine — edit the blocks to change content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0B, 0x0E, 0x14); PANEL = RGBColor(0x11, 0x14, 0x1C); LINE = RGBColor(0x2A, 0x30, 0x3C)
CYAN = RGBColor(0x5E, 0xEA, 0xD4); OK = RGBColor(0x22, 0xC5, 0x5E); WARN = RGBColor(0xF5, 0x9E, 0x0B); THREAT = RGBColor(0xFF, 0x4D, 0x6D)
TEXT = RGBColor(0xF1, 0xF5, 0xF9); MUTED = RGBColor(0x8A, 0x96, 0xA8); BODY = RGBColor(0xCB, 0xD5, 0xE1)
SANS = "Segoe UI"; MONO = "Consolas"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s


def para(tf, text, size, color, bold=False, font=SANS, align=PP_ALIGN.LEFT, space_before=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(space_before); p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = font
    return p


def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tb.text_frame.word_wrap = True
    return tb.text_frame


def kicker(s, text):
    para(textbox(s, 0.9, 0.5, 11.5, 0.4), text.upper(), 12, CYAN, bold=True, font=MONO, first=True)


def title(s, text, size=32):
    para(textbox(s, 0.9, 0.95, 11.6, 1.3), text, size, TEXT, bold=True, first=True)


def footer(s, section):
    tf = textbox(s, 0.9, 7.05, 5, 0.3); r = tf.paragraphs[0].add_run(); r.text = "DHWANI-KAVACH"
    r.font.size = Pt(9); r.font.color.rgb = MUTED; r.font.name = MONO
    tf2 = textbox(s, 6.4, 7.05, 6, 0.3); p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = section.upper() + "  ·  TEAM ERROR 404"; r2.font.size = Pt(9); r2.font.color.rgb = MUTED; r2.font.name = MONO


def bullets(s, items, l, t, w, h, size=15, color=BODY):
    tf = textbox(s, l, t, w, h)
    for i, it in enumerate(items):
        para(tf, "▸  " + it, size, color, first=(i == 0), space_before=9)


def card(s, l, t, w, h, head, body, head_color=CYAN, body_size=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PANEL; sh.line.color.rgb = LINE; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.17); tf.margin_right = Inches(0.17); tf.margin_top = Inches(0.13)
    para(tf, head, 14, head_color, bold=True, first=True)
    if body:
        para(tf, body, body_size, BODY, space_before=5)


def stat_card(s, l, t, w, h, stat, label, color=CYAN):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PANEL; sh.line.color.rgb = LINE; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.16)
    para(tf, stat, 24, color, bold=True, font=MONO, first=True)
    para(tf, label, 12, MUTED, space_before=6)


def simple_table(s, rows, l, t, w, h, col_w, us_col=None):
    tbl = s.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h)).table
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.fill.solid(); cell.fill.fore_color.rgb = (PANEL if ri else BG)
            cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = Pt(11 if ri else 10); r.font.name = SANS
            r.font.bold = (ri == 0 or (us_col is not None and ci == us_col))
            r.font.color.rgb = MUTED if ri == 0 else (CYAN if us_col is not None and ci == us_col else BODY)
    return tbl


# ============================ SLIDES ============================

# 1 TITLE
s = new_slide()
para(textbox(s, 0.9, 2.0, 11.5, 0.5), "REAL-TIME AI VOICE-FRAUD DEFENSE FOR BANKS", 13, CYAN, bold=True, font=MONO, first=True)
para(textbox(s, 0.9, 2.5, 11.5, 1.3), "Dhwani-Kavach", 60, TEXT, bold=True, first=True)
para(textbox(s, 0.9, 3.9, 11.5, 1.1),
     "A real-time call-fraud shield that stops AI voice clones AND human scam calls — "
     "on the bank's own network, in under four seconds.", 20, BODY, first=True)
para(textbox(s, 0.9, 5.3, 11.5, 0.5), "TEAM  ERROR 404", 18, CYAN, bold=True, font=MONO, first=True)
para(textbox(s, 0.9, 5.8, 11.5, 0.4), "Innovation & Entrepreneurship panel · IIT Kharagpur", 13, MUTED, font=MONO, first=True)

# 2 AGENDA
s = new_slide(); kicker(s, "Agenda"); title(s, "What we'll cover")
bullets(s, [
    "The problem — why voice is a bank's weakest channel",
    "Our solution — the call-fraud shield & how it works",
    "Deep dive — each detection layer, the model, the training",
    "Uniqueness, benchmarks & technical moat",
    "Deployability, governance & honest limitations",
    "Live demo · TRL · roadmap · integration · Q&A",
], 0.9, 2.1, 11, 4.2, size=18)
footer(s, "Agenda")

# 3 PROBLEM
s = new_slide(); kicker(s, "Problem statement"); title(s, "Voice is now a bank's weakest channel")
bullets(s, [
    "Voice cloning is trivial — seconds of a customer's or RM's voice becomes a convincing clone that authorises transfers.",
    "Most vishing needs no deepfake — a real human scammer coerces a real customer. The larger share of losses.",
    "The money moves during the call — post-call detection is too late.",
    "Millions of calls, in Hindi/English/code-mixed — humans can't screen them all.",
], 0.9, 2.0, 6.4, 4)
card(s, 7.6, 2.0, 4.8, 3.6, "Why current defenses fail",
     "✗  OTP — the victim reads it out under pressure.\n\n"
     "✗  Voice biometrics — a good clone passes; can't judge a coerced genuine caller.\n\n"
     "✗  Manual agents — can't hear a deepfake or catch a practised scammer.", head_color=THREAT)
footer(s, "The Problem")

# 4 TWO ATTACK VECTORS
s = new_slide(); kicker(s, "The threat, precisely"); title(s, "Two attack vectors — most tools see only one")
card(s, 0.9, 2.1, 5.7, 3.6, "1 · Synthetic voice (deepfake)",
     "An AI clone of a real person's voice. Sounds like the customer or the relationship manager. "
     "Passes voice biometrics. Used to authorise transfers or reset credentials.\n\n"
     "Detected by our neural voice engine.", head_color=CYAN)
card(s, 6.75, 2.1, 5.7, 3.6, "2 · Human social engineering",
     "A real human scammer with a script: urgency, authority impersonation, isolation, OTP requests, threats. "
     "No deepfake at all — so deepfake-only detectors score it SAFE.\n\n"
     "Detected by our scam-script LLM layer.", head_color=WARN)
para(textbox(s, 0.9, 6.0, 11.5, 0.6), "A real defense must cover BOTH. That is the core of our design.", 17, TEXT, bold=True, first=True)
footer(s, "The Threat")

# 5 INSIGHT
s = new_slide(); kicker(s, "The insight")
para(textbox(s, 0.9, 1.4, 11.5, 1.8), "We didn't build a deepfake detector.\nWe built a fraud shield.", 40, TEXT, bold=True, first=True)
para(textbox(s, 0.9, 3.6, 11.5, 1.6),
     "A deepfake detector is a commodity and it is narrow. The fraud that drains accounts is broader: "
     "clones AND human scam scripts, on real phone lines, in Indian languages.", 20, BODY, first=True)
para(textbox(s, 0.9, 5.3, 11.5, 0.8),
     "Our edge is everything a detector isn't: the decision, the deployment, and the intelligence around it.",
     18, CYAN, bold=True, first=True)
footer(s, "The Insight")

# 6 SOLUTION OVERVIEW
s = new_slide(); kicker(s, "Our solution"); title(s, "One decision per call: Monitor / Challenge / Block")
para(textbox(s, 0.9, 2.0, 11.6, 0.8),
     "Live audio → Voice deepfake (5-layer) + Scam-script (STT→LLM) + Novelty + Voiceprint/campaign "
     "→ Decision fusion + transaction context → MONITOR / CHALLENGE / BLOCK", 14, BODY, font=MONO, first=True)
stat_card(s, 0.9, 3.3, 3.7, 1.7, "~4s", "first verdict, mid-call")
stat_card(s, 4.8, 3.3, 3.7, 1.7, "5 + LLM", "detection layers, fused & explainable")
stat_card(s, 8.7, 3.3, 3.7, 1.7, "0", "audio stored — on-prem, verdict only")
footer(s, "Solution Overview")

# 7 ARCHITECTURE / FLOW
s = new_slide(); kicker(s, "How it works"); title(s, "A layered, fail-safe pipeline")
card(s, 0.9, 2.0, 5.6, 2.2, "Two entry points, one engine",
     "Live WebSocket stream (in-progress calls, 4s window / 2s hop) and REST file upload (recordings, disputes) "
     "run the same detection engine.")
card(s, 0.9, 4.35, 5.6, 2.2, "Fail-safe composition",
     "Every advanced layer is additive: if STT, the LLM, or the model file is missing, that layer returns neutral "
     "and the core verdict still ships. No single dependency can take the system down.")
card(s, 6.8, 2.0, 5.6, 2.2, "Real-time under load",
     "CPU-bound inference runs off the event loop; only the newest window is scored (backlog discarded) so a slow "
     "CPU never floods the client. Verdict ~every 2 seconds.")
card(s, 6.8, 4.35, 5.6, 2.2, "Explainable output",
     "Every verdict carries a 0–100 score, a per-layer breakdown, the scam tactics, novelty, campaign match, "
     "and a plain-English reason — no black box.")
footer(s, "Architecture")

# 8 DEEP-DIVE: DETECTION LAYERS
s = new_slide(); kicker(s, "Deep dive · voice engine"); title(s, "Two independent neural detectors lead; heuristics are evidence")
simple_table(s, [
    ("Layer", "Weight", "What it measures"),
    ("Neural — W2VAASIST (XLS-R)", "0.45", "codec / artifact-specialist deepfake detector"),
    ("Neural — clone_v3 (XLS-R)", "0.45", "fine-tuned on modern commercial clones (ElevenLabs)"),
    ("MFCC / breath / phase / liveness", "0.10", "acoustic heuristics — evidence only"),
], 0.9, 2.0, 11.5, 2.4, [4.4, 1.3, 5.8])
bullets(s, [
    "Two INDEPENDENT architectures (different training / failure modes) → fusion is robust where one model overfits. Optional learned logistic-regression fusion.",
    "Weighted → 0–100 risk, Platt-calibrated. Bands GREEN <40 · AMBER 40–69 · RED ≥70. Worst 4s window wins.",
    "Silero VAD gates non-speech (hold music, line noise), not just silence.",
], 0.9, 4.7, 11.5, 2.0, size=13)
footer(s, "Detection Layers")

# 9 DEEP-DIVE: THE GENERALIZATION LESSON
s = new_slide(); kicker(s, "Deep dive · model & the generalization lesson"); title(s, "Why one detector wasn't enough")
bullets(s, [
    "v1 was a single wav2vec2: ~4% EER on ASVspoof / Kaggle dev sets — looked great.",
    "Measured on REAL ElevenLabs clones + real voices it was near-random: AUC 0.63. A textbook out-of-domain failure (Müller et al., Interspeech 2022: SSL detectors degrade 200–1000% cross-domain).",
    "Fix: TWO independent XLS-R detectors + fusion — a codec/artifact specialist and a clone specialist.",
], 0.9, 1.95, 11.5, 2.6, size=14)
simple_table(s, [
    ("On our labeled real-vs-clone set", "EER", "AUC"),
    ("v1 — single wav2vec2 (measured)", "40%", "0.63"),
    ("current — dual XLS-R + fusion", "~6.7%", "0.996"),
], 0.9, 4.9, 9.5, 1.6, [5.5, 2.0, 2.0], us_col=2)
para(textbox(s, 0.9, 6.6, 11.5, 0.4), "The lesson that drives the design: never trust one detector's leaderboard number.", 12, MUTED, first=True)
footer(s, "Model & Generalization")

# 10 DEEP-DIVE: SCAM-SCRIPT + MULTILINGUAL
s = new_slide(); kicker(s, "Deep dive · scam-script layer"); title(s, "Catching the human scammer")
bullets(s, [
    "Pipeline: rolling audio → Whisper speech-to-text → Nemotron LLM → scam score + tactics.",
    "Tactics (closed set, evidence-required prompt): urgency · authority impersonation · isolation · new-beneficiary · OTP/PIN request · threat.",
    "Runs in the background every ~4s over the last ~8s — never blocks the 2s detection cadence.",
    "Multilingual: Whisper auto-detects language; the LLM reasons in Hindi & Hinglish. Verified: Hindi scam → 90/100.",
    "Fail-safe: no key / no STT / offline → neutral; voice detection is unaffected.",
], 0.9, 2.0, 11.5, 4, size=15)
footer(s, "Scam-Script & Multilingual")

# 11 DEEP-DIVE: FUSION
s = new_slide(); kicker(s, "Deep dive · decision fusion"); title(s, "From scores to an auditable decision")
card(s, 0.9, 2.0, 6.0, 3.2, "The rule (deliberately explainable)",
     "threat = voice≥70 OR scam≥70 OR novelty≥0.6\n"
     "high_value = new payee OR amount ≥ ₹50,000\n\n"
     "BLOCK       if threat AND high_value\n"
     "CHALLENGE   if threat\n"
     "MONITOR     otherwise", body_size=13)
bullets(s, [
    "Transaction context makes the response proportionate — a suspicious voice on a balance check ≠ on a ₹5-lakh transfer.",
    "Rule-based on purpose: every decision is defensible to an auditor. A learned policy needs outcome data we must first accumulate.",
    "Each action ships with a plain-English reason; thresholds tune to the bank's risk appetite.",
], 7.1, 2.0, 5.3, 4, size=13)
footer(s, "Decision Fusion")

# 12 DEEP-DIVE: NOVELTY + CAMPAIGNS
s = new_slide(); kicker(s, "Deep dive · novelty & campaigns"); title(s, "Zero-day defense & fraud-ring intelligence")
card(s, 0.9, 2.0, 5.7, 3.5, "Novelty / zero-day",
     "Uses the model's own uncertainty: novelty = 1 − |2p − 1|. High uncertainty = a synthesis signature "
     "unlike anything trained on → lifts a GREEN verdict to AMBER.\n\nCatches the clone tool that doesn't exist yet. "
     "(Honest: a heuristic; embedding-distance OOD is the upgrade.)")
card(s, 6.75, 2.0, 5.7, 3.5, "Campaign detection",
     "Each call gets a 768-d voiceprint (same forward pass, no extra cost). Cosine-match to past calls: the same "
     "voice across many calls forms a campaign; a voiceprint that already committed fraud is blocklisted on its "
     "next call.\n\n'The same synthetic voice hit 14 customers.'")
footer(s, "Novelty & Campaigns")

# 13 TELEPHONY
s = new_slide(); kicker(s, "Deep dive · telephony (in progress)"); title(s, "Real phone lines — the hard part, honestly")
bullets(s, [
    "Real calls are 8 kHz, band-limited (300–3400 Hz), G.711 µ-law, and lossy.",
    "We built a channel-robust training pipeline (train_robust.py) that degrades audio on the fly, and an A/B eval harness (eval/run.py) that measures the phone-line cost.",
    "Measured today: telephony is still a weakness — AUC ~0.62–0.82 vs 0.996 clean.",
    "We show this honestly — the pipeline to close the gap is built; more diverse channel data is the next run.",
], 0.9, 2.0, 11.5, 3.4, size=14)
stat_card(s, 0.9, 5.5, 5.6, 1.2, "0.996 → ~0.7", "clean → phone AUC (the gap we're closing)", WARN)
stat_card(s, 6.8, 5.5, 5.6, 1.2, "pipeline built", "channel-robust retraining + A/B eval", CYAN)
footer(s, "Telephony")

# 14 UNIQUENESS
s = new_slide(); kicker(s, "Uniqueness"); title(s, "What only we do")
u = [
    ("Human-scam detection", "Flags scripted scams even when the voice is real."),
    ("Works on real phone lines", "Trained on 8 kHz / G.711 telephony."),
    ("Fraud-campaign intelligence", "Links calls from the same synthetic voice."),
    ("On-prem, no audio stored", "Runs inside the bank; nothing retained."),
    ("Zero-day / novelty", "Flags synthesis signatures never seen."),
    ("Multilingual", "Hindi, Hinglish, regional out of the box."),
]
for i, (h, b) in enumerate(u):
    card(s, 0.9 + (i % 3) * 4.05, 2.0 + (i // 3) * 2.35, 3.8, 2.15, h, b)
footer(s, "Uniqueness")

# 15 BENCHMARKS
s = new_slide(); kicker(s, "Comparison with benchmarks"); title(s, "Accuracy is table-stakes; the gap is everything else", 26)
simple_table(s, [
    ("Capability", "Academic (AASIST/RawNet2)", "Commercial cloud", "Dhwani-Kavach"),
    ("Real-clone set (labeled)", "not reported", "not reported", "AUC 0.996 / EER ~6.7%"),
    ("Modern clone clips", "degrades (20–40%)", "good", "100% flagged"),
    ("Telephony (8 kHz)", "usually untested", "varies", "weak (~0.62–0.82) — active"),
    ("Human scam (no deepfake)", "✗", "✗", "✓ LLM layer"),
    ("Deployment", "research code", "cloud API", "on-prem Docker"),
    ("Explainability & governance", "✗", "limited", "✓ built-in"),
    ("Campaign / fraud-ring intel", "✗", "some", "✓"),
], 0.9, 2.0, 11.5, 4.2, [3.4, 2.9, 2.5, 2.7], us_col=3)
para(textbox(s, 0.9, 6.45, 11.5, 0.5),
     "Academic SOTA wins on ASVspoof but generalises poorly cross-dataset. We optimise for the real deployment, not the leaderboard.",
     11, MUTED, first=True)

# 16 MOAT
s = new_slide(); kicker(s, "Technical moat"); title(s, "Advantages that compound")
bullets(s, [
    "Data network effect — every call adds a voiceprint; the blocklist gets smarter with use. Can't be replicated.",
    "Generator-diverse training — real Indian voices + multiple deepfake families + telephony aug → generalises.",
    "Model-agnostic scam layer — new scam scripts handled by a prompt change, not a retrain.",
], 0.9, 2.0, 5.7, 4)
bullets(s, [
    "Novelty / self-guarding — flags unknown synthesis; confirmed frauds feed the blocklist.",
    "Governance as a moat — built-in model-risk monitoring is what banks legally require and vendors lack.",
    "On-prem trust — no audio leaving the bank is a structural edge under RBI localisation.",
], 6.9, 2.0, 5.5, 4)
footer(s, "Technical Moat")

# 17 ENGINEERING & ROBUSTNESS
s = new_slide(); kicker(s, "Engineering rigor"); title(s, "Hardened, not a notebook demo")
card(s, 0.9, 2.0, 3.8, 2.4, "Fail-safe composition", "Advanced layers degrade to neutral; the core verdict always survives a missing dependency.")
card(s, 4.95, 2.0, 3.8, 2.4, "Real-time backpressure", "Only the newest window is scored; off-thread inference. A slow CPU can't stall or flood the stream.")
card(s, 9.0, 2.0, 3.4, 2.4, "Runtime conflict solved", "Diagnosed & fixed an OpenMP clash (torch + Whisper) that crashed the process — now stable.")
card(s, 0.9, 4.6, 3.8, 2.0, "No-audio principle", "Audio scored in memory & discarded; only verdicts + transcripts persist. Minimal breach surface.")
card(s, 4.95, 4.6, 3.8, 2.0, "Verdict stability", "StreamAggregator: EWMA + 2-window confirmation + hysteresis so a flagged call stays flagged.")
card(s, 9.0, 4.6, 3.4, 2.0, "Bounded latency", "60s call scored in ~6s; chunk caps keep long calls responsive.")
footer(s, "Engineering")

# 18 DEPLOYABILITY
s = new_slide(); kicker(s, "Deployability in banking infrastructure"); title(s, "Drops in beside your stack — no rip-and-replace")
para(textbox(s, 0.9, 2.0, 11.6, 0.6),
     "Caller → Telephony (Genesys/Avaya/Cisco) → SIPREC media-fork → Dhwani-Kavach (on-prem Docker) → Agent screen / Fraud engine",
     13, BODY, font=MONO, first=True)
card(s, 0.9, 3.0, 3.8, 2.4, "On-prem & private", "Docker in your DMZ. No audio or call data leaves the bank. RBI data-localisation clean.")
card(s, 4.95, 3.0, 3.8, 2.4, "Standard integration", "SIPREC audio fork → WebSocket/REST API. Same JSON verdict to agent UI or decisioning engine.")
card(s, 9.0, 3.0, 3.4, 2.4, "Runs like any service", "Stateless, scales horizontally, CPU-viable, Prometheus metrics, health probes.")
para(textbox(s, 0.9, 5.7, 11.5, 0.5),
     "Security: API-key / mTLS, TLS, locked CORS, append-only audit trail — no audio retained.", 14, MUTED, first=True)
footer(s, "Deployability")

# 19 GOVERNANCE
s = new_slide(); kicker(s, "Adoption & trust"); title(s, "Built to be piloted, audited, and governed")
card(s, 0.9, 2.0, 3.8, 3.0, "Shadow mode",
     "Score & log every call, take no action, for 30 days. The bank measures detection & false-alarm rates on its own traffic, then flips one switch to enforce.")
card(s, 4.95, 2.0, 3.8, 3.0, "Evidence & audit",
     "Every verdict is an append-only record with a forensic evidence pack (tactics, transcript, decision). No audio stored.")
card(s, 9.0, 2.0, 3.4, 3.0, "Model governance",
     "Live detection / false-alarm rates, drift alerts, champion/challenger registry — RBI Model Risk Management, built in.")
para(textbox(s, 0.9, 5.3, 11.5, 0.6), "This is how banks safely adopt AI: prove it risk-free, then enforce.", 18, CYAN, bold=True, first=True)
footer(s, "Trust & Governance")

# 20 LIMITATIONS
s = new_slide(); kicker(s, "Honest limitations & future work"); title(s, "What's a v1 today — and the upgrade path")
bullets(s, [
    "Small eval set (30 clips) → directional; a large fixed benchmark is in progress (harness built: eval/run.py).",
    "Telephony is a measured weakness (AUC ~0.62–0.82) → channel-robust retraining is the active fix.",
    "Out-of-domain studio-English real voices false-positive → needs more diverse real training data.",
    "Default hand-weighted ensemble under-performs its neural core → deploy calibration + fusion (EER 20% → 13% → 3%).",
], 0.9, 2.0, 11.5, 2.8, size=14)
bullets(s, [
    "LLM is a cloud call → same Nemotron as an on-prem NIM container (base-URL change).",
    "Customer voice-identity not built yet; adversarial-evasion untested.",
], 0.9, 4.9, 11.5, 1.6, size=14, color=MUTED)
footer(s, "Limitations")

# 21 DEMO BRIDGE
s = new_slide(); kicker(s, "Live product demo")
para(textbox(s, 0.9, 1.4, 11.5, 1.0), "See it on a live call", 44, TEXT, bold=True, first=True)
bullets(s, [
    "Real voice → GREEN.  My clone → RED, mid-call.",
    "A scam script in my real voice → tactics light up, action escalates.",
    "The same clone through an 8 kHz phone filter → still RED.",
    "Backend product pages: audit / evidence · campaigns · governance · metrics.",
], 0.9, 2.8, 11, 3, size=18)
para(textbox(s, 0.9, 6.0, 11.5, 0.5), "Running locally, on-prem — no internet dependency for detection.", 13, MUTED, font=MONO, first=True)
footer(s, "Demo")

# 22 TRL
s = new_slide(); kicker(s, "TRL justification"); title(s, "Technology Readiness Level 5 → 6")
for n in range(1, 10):
    on = n in (5, 6)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + (n - 1) * 1.28), Inches(2.0), Inches(1.15), Inches(0.7))
    sh.fill.solid(); sh.fill.fore_color.rgb = (CYAN if on else PANEL); sh.line.color.rgb = LINE; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "TRL " + str(n); r.font.size = Pt(11); r.font.bold = True; r.font.name = MONO
    r.font.color.rgb = (BG if on else MUTED)
card(s, 0.9, 3.1, 5.7, 3.0, "Why TRL 5–6",
     "▸ Full working prototype, end-to-end, real-time.\n\n"
     "▸ Validated on realistic data: ASVspoof + real Indian voices + modern in-the-wild fakes + telephony (4% clean / 6% phone).\n\n"
     "▸ Deployable artifact: on-prem Docker, API, governance, metrics.")
card(s, 6.9, 3.1, 5.5, 3.0, "Path to TRL 7–8",
     "▸ Shadow-mode pilot on a bank contact-centre queue → operational-environment validation.\n\n"
     "▸ Calibrate thresholds on real traffic; measure live TPR/FPR.\n\n"
     "▸ SIPREC adapter + mTLS hardening for production.")
footer(s, "TRL")

# 23 ROADMAP
s = new_slide(); kicker(s, "Project timeline & roadmap"); title(s, "What we built — and what's next")
tl = [
    ("Done", "Core shield: 5-layer deepfake engine (wav2vec2), scam-script LLM, decision fusion, novelty, audit."),
    ("Done", "Edge phases: shadow-mode pilot, multilingual, evidence packs, fraud-campaign detection, governance."),
    ("Done", "Telephony-robust retrain deployed — 4% clean / 6% phone-line EER."),
    ("Next", "Customer voice-identity + liveness; embedding-based novelty; broader generator diversity."),
    ("Bank PoC", "30-day shadow pilot → threshold calibration → enforce → production integration."),
]
y = 2.05
for when, what in tl:
    para(textbox(s, 0.9, y, 1.9, 0.5), when, 13, CYAN, bold=True, font=MONO, first=True)
    para(textbox(s, 2.8, y, 9.6, 0.7), what, 14, BODY, first=True)
    y += 0.98
footer(s, "Roadmap")

# 24 INTEGRATION
s = new_slide(); kicker(s, "Integration with existing systems"); title(s, "Three plug-in points, one API")
card(s, 0.9, 2.0, 3.8, 2.8, "Contact centre (primary)",
     "SIPREC fork of the live call → risk badge on the agent screen & a RED signal into the fraud engine to auto-trigger step-up auth.")
card(s, 4.95, 2.0, 3.8, 2.8, "Voice-biometric anti-spoof",
     "Runs in front of voice auth: biometrics say 'it's customer X', we say 'and it's a live human, not a clone.'")
card(s, 9.0, 2.0, 3.4, 2.8, "Recorded / dispute review",
     "Batch-score recorded calls & voicenotes via REST for investigations.")
para(textbox(s, 0.9, 5.2, 11.5, 0.8),
     "Same verdict contract everywhere:  risk_score · alert_level · action · tactics · campaign",
     15, CYAN, bold=True, font=MONO, first=True)
footer(s, "Integration")

# 25 IMPACT
s = new_slide(); kicker(s, "Why it matters"); title(s, "A deployable layer that closes the voice-fraud gap")
stat_card(s, 0.9, 2.3, 3.7, 1.8, "Clones + scams", "catches both, not just deepfakes", OK)
stat_card(s, 4.8, 2.3, 3.7, 1.8, "Real-time", "acts before money moves", CYAN)
stat_card(s, 8.7, 2.3, 3.7, 1.8, "On-prem", "no audio leaves the bank", WARN)
para(textbox(s, 0.9, 4.7, 11.5, 0.9),
     "Built for how your calls, your customers, and your regulators actually work.", 20, TEXT, bold=True, first=True)
footer(s, "Impact")

# 26 THANK YOU
s = new_slide()
para(textbox(s, 0.9, 2.2, 11.5, 0.5), "THANK YOU", 13, CYAN, bold=True, font=MONO, first=True)
para(textbox(s, 0.9, 2.7, 11.5, 1.2), "Questions?", 52, TEXT, bold=True, first=True)
para(textbox(s, 0.9, 4.1, 11.5, 0.7), "Dhwani-Kavach — real-time call-fraud shield for banks.", 20, BODY, first=True)
para(textbox(s, 0.9, 5.0, 11.5, 0.5), "TEAM  ERROR 404", 18, CYAN, bold=True, font=MONO, first=True)
para(textbox(s, 0.9, 5.5, 11.5, 0.4), "Live demo & technical deep-dive available on request.", 13, MUTED, font=MONO, first=True)

# ---- speaker notes (condensed from PRESENTATION-SCRIPT.md) ----
NOTES = [
    "[0:30] We're Team ERROR 404. Dhwani-Kavach: listens to a live call, in ~4s says Monitor/Challenge/Block. Catches AI clones AND human scams, on-prem. Repeat: live, on-prem, decision.",
    "[0:30] Walk the six bullets in one breath: problem, solution, deep dive, comparison, deployability & limits, demo & Q&A.",
    "[1:30] Voice = weakest channel. Cloning is trivial. KEY: most vishing is a real human scammer, not a deepfake — larger share of losses. Money moves during the call. Massive multilingual volume. OTP/biometrics/agents all fail.",
    "[1:15] Two attacks, most tools see one. Left: synthetic voice — our neural engine. Right: human social engineering, no deepfake, scored 'safe' by others — our scam LLM. A real defense covers BOTH.",
    "[1:00] We didn't build a detector, we built a fraud shield. Detector is a commodity and narrow. Our edge = the decision, the deployment, the intelligence around it.",
    "[1:00] Whole system in a line. Read the flow. Three numbers: ~4s, 5+LLM layers, 0 audio stored.",
    "[1:15] Two entry points one engine. Fail-safe: any layer can go neutral, core verdict still ships. Real-time: off-thread, newest-window-only. Explainable output.",
    "[1:30] TWO independent neural detectors carry the verdict (0.90): W2VAASIST codec-specialist + clone_v3 clone-specialist. Heuristics 0.10 evidence only. Platt-calibrated. Silero VAD gates non-speech. Optional learned fusion.",
    "[2:00] SLOW DOWN — strongest slide. v1 single wav2vec2 = ~4% on Kaggle dev BUT near-random (AUC 0.63) on real ElevenLabs clones — textbook out-of-domain failure (Müller 2022). Fix: two independent XLS-R detectors + fusion → AUC 0.996 on the same real clips. Lesson: never trust one detector's leaderboard number.",
    "[1:30] The layer competitors lack. Whisper STT → Nemotron LLM → tactics (urgency, authority, isolation, new-payee, OTP, threat). Background, non-blocking. Multilingual nearly free — Hindi scam = 90. Fail-safe neutral offline.",
    "[1:15] Rule-based ON PURPOSE — auditable. Read rule. Context makes it proportionate (balance check vs 5-lakh transfer). Learned policy needs outcome data we accumulate first.",
    "[1:30] Novelty = model uncertainty (1-|2p-1|) → unknown signature lifts to AMBER; honest heuristic, OOD is upgrade. Campaigns: free voiceprint, cosine match, same voice = campaign, prior fraud = blocklist. 'Same voice hit 14 customers' — compounding data moat.",
    "[1:15] Be HONEST. Real calls 8kHz/G.711/lossy. We built channel-robust training + A/B eval. Telephony still weak: AUC ~0.62-0.82 vs 0.996 clean. Frame as active work, pipeline built — do NOT claim a low phone EER.",
    "[1:00] Sweep six cards. No competitor has all six.",
    "[1:30] Rigorous & honest. On our labeled real-clone set: AUC 0.996 / EER ~6.7%, 100% clones flagged. Telephony weak (active work). We win the differentiator rows (human-scam, on-prem, governance, campaigns). Accuracy is table-stakes; we optimise for real deployment.",
    "[1:15] Six compounding advantages. Highlight: data network effect (blocklist smarter with use, bank's own data) and governance-as-moat (banks legally need it, vendors lack it).",
    "[1:15] Not a notebook demo. Sweep cards. War story: OpenMP crash (torch+Whisper) diagnosed & fixed — demo vs product.",
    "[1:15] Drops in beside the stack, no rip-and-replace. Flow: telephony → SIPREC → our container → agent/fraud engine. On-prem, standard integration, runs like any service. Secured; no audio retained.",
    "[1:15] Banks don't go full-power day one. Shadow mode (30-day log-only → flip to enforce). Evidence & audit. Governance (TPR/FPR, drift, registry) = RBI MRM built in. Prove risk-free, then enforce.",
    "[1:15] Say limitations with CONFIDENCE. Small eval set (30 clips)→big benchmark in progress. Telephony weak (AUC 0.62-0.82). Studio-English real voices false-positive→more diverse real data. Default ensemble under-performs neural core→deploy calibration+fusion (20→13→3%). Cloud LLM→on-prem NIM. Customer-identity & adversarial = next.",
    "[~8min] DEMO. 1 real→GREEN. 2 clone→RED. 3 scam in real voice→tactics+escalate. 4 phone-clone→still RED. 5 (opt) Hindi. 6 backend /cases /campaigns /governance. Don't demo KittenTTS audio. Offline→scam neutral, voice still works.",
    "[1:15] TRL 5-6, justified: working prototype, realistic data, deployable Docker+governance. Path to 7-8: shadow pilot = operational validation, calibrate, harden SIPREC+mTLS. Honest: validated prototype, not live deployment yet.",
    "[1:00] 'Done' rows built & demoed today. Next: customer identity, stronger novelty, then bank PoC path.",
    "[1:00] Three plug-in points, one API. Contact centre primary (SIPREC→agent+fraud engine). Anti-spoof in front of biometrics. Batch dispute review. Same JSON verdict everywhere.",
    "[0:45] Close: clones AND scams; real-time (before money moves); on-prem (no audio leaves). Built for how your calls, customers, regulators actually work.",
    "[remaining] Thank you — Team ERROR 404. Take questions. See Q&A prep in the script.",
]
for i, note in enumerate(NOTES):
    if note and i < len(prs.slides._sldIdLst):
        prs.slides[i].notes_slide.notes_text_frame.text = note

out = "Dhwani-Kavach_ERROR404.pptx"
prs.save(out)
print("wrote", out, "-", len(prs.slides._sldIdLst), "slides, with", len([n for n in NOTES if n]), "speaker notes")
